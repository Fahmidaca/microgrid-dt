"""
Generate microgrid_paper_progress.pptx  - a slide deck summarising all
CSE-side work on the microgrid digital-twin paper so far.

Sections:
    1. Title + agenda
    2. UCI Grid Stability pipeline (yesterday)
    3. Synthetic dataset + compliance + forecasting (today)
    4. EEE-side simulation summary
    5. Honesty notes, gaps, and target venues

Uses python-pptx. Regenerate any time by running:
    python src/generate_presentation.py
Output: presentation/microgrid_paper_progress.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"; OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR = ROOT / "figures"

# =========================================================================
#  THEME
# =========================================================================
DARK = RGBColor(0x26, 0x46, 0x53)      # #264653 - navy background
TEAL = RGBColor(0x2A, 0x9D, 0x8F)      # #2a9d8f - accent
CORAL = RGBColor(0xE7, 0x6F, 0x51)     # #e76f51 - accent 2
CREAM = RGBColor(0xE9, 0xC4, 0x6A)     # #e9c46a - highlight
LIGHT = RGBColor(0xF4, 0xF1, 0xDE)     # off-white text on dark
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x60, 0x60, 0x60)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)   # 16:9 widescreen


# =========================================================================
#  HELPERS
# =========================================================================
def _paint_bg(slide, color=DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_text(slide, text, left, top, width, height, size=20,
              bold=False, color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb


def _add_bullets(slide, bullets, left, top, width, height, size=16,
                 color=LIGHT, spacing_after=8, bullet_char="•"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing_after)
        run = p.add_run()
        run.text = f"{bullet_char}  {b}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
    return tb


def _add_line(slide, x1, y1, x2, y2, color=TEAL, weight=2.0):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def _add_image(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), left, top,
                                        width=width, height=height)
    return None


def _add_footer(slide, label, page):
    _add_text(slide, label, Inches(0.4), Inches(7.05), Inches(9),
              Inches(0.35), size=10, color=GRAY)
    _add_text(slide, f"{page}", Inches(12.5), Inches(7.05),
              Inches(0.5), Inches(0.35), size=10, color=GRAY,
              align=PP_ALIGN.RIGHT)


def _add_table(slide, headers, rows, left, top, width, height,
               header_color=TEAL, body_size=13, header_size=14):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top,
                                       width, height)
    tbl = tbl_shape.table
    # header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_color
        cell.text_frame.text = ""
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = h
        run.font.size = Pt(header_size); run.font.bold = True
        run.font.color.rgb = WHITE; run.font.name = "Calibri"
    # rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (WHITE if i % 2 == 1
                                        else RGBColor(0xF0, 0xF0, 0xF0))
            cell.text_frame.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(body_size)
            run.font.color.rgb = DARK
            run.font.name = "Calibri"
    return tbl_shape


# =========================================================================
#  SLIDE BUILDERS
# =========================================================================
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    # coral accent bar spans the whole text stack
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(1.7), Inches(0.15), Inches(3.9))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    _add_text(s, "Microgrid Digital Twin",
              Inches(1.0), Inches(1.7), Inches(11.5), Inches(1.2),
              size=54, bold=True, color=WHITE)
    _add_text(s, "Stability, Power Quality, and Forecasting",
              Inches(1.0), Inches(2.95), Inches(11.5), Inches(0.9),
              size=34, color=TEAL)
    _add_text(s,
              "Progress report  -  CSE side ML pipelines  +  "
              "EEE side simulation",
              Inches(1.0), Inches(3.95), Inches(11.5), Inches(0.7),
              size=22, bold=True, color=CREAM)
    _add_text(s,
              "Target venues:  IEEE Trans. Smart Grid (Q1, IF~10)  |  "
              "Applied Energy (Q1, IF~11)",
              Inches(1.0), Inches(4.9), Inches(11.5), Inches(0.6),
              size=19, color=LIGHT)
    _add_text(s, "github.com/Fahmidaca/microgrid-dt-1",
              Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.5),
              size=17, color=CREAM, font="Consolas")


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Agenda", Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=32, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)
    agenda = [
        "The problem: renewable microgrid stability and IEEE 519 compliance",
        "Project layout: two halves (CSE + EEE) that meet in the middle",
        "Part 1  -  UCI Grid Stability pipeline (public benchmark)",
        "Part 2  -  Synthetic dataset + compliance classifier + forecasting",
        "Part 3  -  EEE side physics simulation",
        "How the two halves meet: the AI-Optimized Control link",
        "Honesty notes, gaps, and next steps",
    ]
    _add_bullets(s, agenda, Inches(1.0), Inches(1.7), Inches(11),
                 Inches(5.0), size=20, color=DARK, spacing_after=14)
    _add_footer(s, "Agenda", 2)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "The problem", Inches(0.7), Inches(0.5),
              Inches(12), Inches(0.7),
              size=32, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)
    _add_text(s,
              "Renewable microgrids are hard to keep stable and IEEE 519 compliant "
              "as PV / wind penetration grows.",
              Inches(0.7), Inches(1.6), Inches(12), Inches(1.0),
              size=20, color=DARK)
    q = [
        "How stable is the grid at any moment?",
        "How much can we trust the ML model's stability call?",
        "Will voltage-THD breach the IEEE 519 5 % limit in the next 5-30 minutes?",
        "Do renewable-heavy operating scenarios pass IEEE 519 without help?",
        "Can an AI controller do better than fixed control laws?",
    ]
    _add_text(s, "The five questions this project answers:",
              Inches(0.7), Inches(2.8), Inches(12), Inches(0.5),
              size=18, bold=True, color=TEAL)
    _add_bullets(s, q, Inches(1.0), Inches(3.3), Inches(11.5),
                 Inches(3.5), size=18, color=DARK, spacing_after=10)
    _add_footer(s, "Motivation", 3)


def slide_two_halves(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Two halves that meet in the middle",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=32, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    # LEFT box: CSE
    left = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.8), Inches(1.7),
                              Inches(5.7), Inches(4.4))
    left.fill.solid(); left.fill.fore_color.rgb = TEAL
    left.line.color.rgb = TEAL
    tf = left.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "CSE  /  ML SIDE"
    run.font.size = Pt(20); run.font.bold = True; run.font.color.rgb = WHITE
    for txt in [
        "",
        "UCI Grid Stability benchmark",
        "5 classifiers x 5 seeds",
        "Best: MLP  96.05 % accuracy",
        "",
        "NEW METRIC:",
        "tau-robustness margin",
        "(no prior paper reports it)",
        "",
        "Synthetic 50k-row dataset",
        "3-class IEEE 519 compliance classifier",
        "Multi-horizon THD forecasting",
    ]:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = txt
        r.font.size = Pt(14); r.font.color.rgb = WHITE

    # RIGHT box: EEE
    right = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(6.9), Inches(1.7),
                               Inches(5.7), Inches(4.4))
    right.fill.solid(); right.fill.fore_color.rgb = CORAL
    right.line.color.rgb = CORAL
    tf = right.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = "EEE  /  PHYSICS SIDE"
    run.font.size = Pt(20); run.font.bold = True; run.font.color.rgb = WHITE
    for txt in [
        "",
        "Bangladesh PV microgrid (50 Hz, 230 V)",
        "6-pulse rectifier + nonlinear loads",
        "",
        "Simulink builder + Python port",
        "SRF Active Power Filter",
        "FFT-based THD analyzer",
        "Battery ageing + BDT cost model",
        "",
        "IEEE 14-bus parametric study",
        "10 scenarios, IEEE 519 verdict",
        "(EEE side is currently completing this)",
    ]:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = txt
        r.font.size = Pt(14); r.font.color.rgb = WHITE

    # MEETING POINT
    meet = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(3.3), Inches(6.2),
                              Inches(6.7), Inches(0.7))
    meet.fill.solid(); meet.fill.fore_color.rgb = DARK
    meet.line.color.rgb = DARK
    tf = meet.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = ("They meet at the AI-Optimized Control scenario  "
                "(the CSE ML model becomes the controller)")
    run.font.size = Pt(13); run.font.bold = True
    run.font.color.rgb = WHITE
    _add_footer(s, "Project layout", 4)


def slide_uci_dataset(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Part 1  -  UCI Grid Stability dataset",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=30, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)
    _add_text(s,
              "Schafer et al. 2016, Eur. Phys. J. Special Topics 225, 569",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.4),
              size=14, color=TEAL)
    _add_bullets(s, [
        "10,000 simulated 4-bus decentralised smart grid snapshots",
        "12 features:  tau1..4 (reaction times), p1..4 (powers), g1..4 (elasticities)",
        "Binary label:  stable / unstable  (64 % unstable, 36 % stable)",
        "Peer-reviewed, still cited in 2024-25 papers - benchmarks don't expire, methods do",
    ], Inches(0.8), Inches(2.0), Inches(12), Inches(2.5),
       size=17, color=DARK, spacing_after=10)

    # Class distribution image
    _add_image(s, FIG_DIR / "eda" / "01_class_distribution.png",
               Inches(0.8), Inches(4.4), width=Inches(4.5))
    _add_image(s, FIG_DIR / "eda" / "04_tau_stability_boundary.png",
               Inches(5.6), Inches(4.4), width=Inches(7.2))
    _add_footer(s, "UCI dataset", 5)


def slide_uci_methodology(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Part 1  -  Methodology",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=30, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)
    _add_text(s,
              "5 classifiers  x  5 seeds  =  25 independent training runs",
              Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
              size=20, bold=True, color=TEAL)
    _add_bullets(s, [
        "Logistic Regression, Random Forest, HistGB, XGBoost, MLP",
        "80 / 20 stratified train-test split, deterministic seeding",
        "Multi-seed setup = we report mean +/- std, not a single lucky run",
    ], Inches(0.8), Inches(2.2), Inches(12), Inches(1.6),
       size=16, color=DARK, spacing_after=8)

    _add_text(s, "Statistical rigor",
              Inches(0.7), Inches(4.0), Inches(12), Inches(0.5),
              size=20, bold=True, color=TEAL)
    _add_bullets(s, [
        "Paired McNemar test for every pair of models  (10 pairs)",
        "Holm-Bonferroni correction for multiple comparisons",
        "Bootstrap 95 % confidence intervals (2000 iterations)",
        "Result: all 10 pairwise differences significant at p < 0.05",
    ], Inches(0.8), Inches(4.6), Inches(12), Inches(2.2),
       size=16, color=DARK, spacing_after=8)
    _add_footer(s, "UCI methodology", 6)


def slide_uci_results(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Part 1  -  Results (test accuracy, mean +/- std over 5 seeds)",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=26, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_table(s,
               headers=["Model", "Accuracy", "95 % CI", "Macro-F1"],
               rows=[
                   ["MLP (winner)",       "96.05 % +/- 0.62 %", "[95.65, 96.43]", "95.72 %"],
                   ["HistGB",              "94.72 %",             "[94.28, 95.13]", "94.24 %"],
                   ["XGBoost",             "94.39 %",             "[93.92, 94.83]", "93.87 %"],
                   ["Random Forest",       "92.16 %",             "[91.64, 92.70]", "91.38 %"],
                   ["Logistic Regression", "81.50 %",             "[80.72, 82.29]", "79.58 %"],
               ],
               left=Inches(0.8), top=Inches(1.7),
               width=Inches(11.7), height=Inches(2.5),
               header_color=TEAL, body_size=15, header_size=15)

    _add_text(s,
              "All 10 pairwise differences statistically significant "
              "(McNemar + Holm-Bonferroni, p < 0.05)",
              Inches(0.7), Inches(4.4), Inches(12), Inches(0.5),
              size=15, color=GRAY, align=PP_ALIGN.CENTER)

    _add_text(s,
              "The MLP wins on accuracy - but that is only half the story.",
              Inches(0.7), Inches(5.1), Inches(12), Inches(0.5),
              size=17, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    _add_text(s, "-> the other half is on the next slide",
              Inches(0.7), Inches(5.7), Inches(12), Inches(0.5),
              size=13, color=GRAY, align=PP_ALIGN.CENTER)
    _add_footer(s, "UCI results", 7)


def slide_tau_margin(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _add_text(s, "Novel contribution:  tau-robustness margin",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=WHITE)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CREAM, weight=2.5)

    _add_text(s, "The gap in prior work",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
              size=18, bold=True, color=CREAM)
    _add_text(s,
              "Every existing paper on this benchmark reports only "
              "classification accuracy.  But a grid operator needs to "
              "know how much they can trust the model's decision.",
              Inches(0.8), Inches(2.0), Inches(12), Inches(1.2),
              size=15, color=LIGHT)

    _add_text(s, "What we defined",
              Inches(0.7), Inches(3.4), Inches(12), Inches(0.5),
              size=18, bold=True, color=CREAM)
    _add_bullets(s, [
        "For each correctly-classified test point, find the smallest tau-perturbation "
        "that flips the model's prediction",
        "Binary search over 16 sign patterns on tau1..tau4  (L_inf norm)",
        "Clipped to the simulation range [0.5, 10] seconds",
        "Result: a per-sample safety buffer in seconds",
    ], Inches(0.8), Inches(3.9), Inches(12), Inches(2.5),
       size=14, color=LIGHT, spacing_after=6)

    _add_text(s,
              "Larger margin  =  the classifier's decision survives more "
              "consumer reaction-time degradation before it flips",
              Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
              size=14, color=CREAM, align=PP_ALIGN.CENTER)
    _add_footer(s, "Novelty", 8)


def slide_tau_insight(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Key insight:  accuracy  !=  robustness",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_image(s, FIG_DIR / "robustness" / "02_margin_per_class.png",
               Inches(0.6), Inches(1.6), width=Inches(6.5))

    # right-hand insight box
    _add_text(s, "What the numbers show",
              Inches(7.3), Inches(1.6), Inches(5.7), Inches(0.5),
              size=18, bold=True, color=CORAL)
    _add_bullets(s, [
        "MLP  =  best accuracy  (96.05 %)  but tightest tau-margin",
        "LogReg  =  lowest accuracy  (81.5 %)  but LARGEST margin",
        "So 'more accurate' does not mean 'more deployable'",
        "Prior accuracy-only reports hide this trade-off completely",
    ], Inches(7.3), Inches(2.2), Inches(5.7), Inches(3.2),
       size=14, color=DARK, spacing_after=10)

    _add_text(s,
              "This is the publishable insight of the CSE-side work.",
              Inches(0.7), Inches(6.6), Inches(12), Inches(0.5),
              size=16, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    _add_footer(s, "Novelty - insight", 9)


def slide_synth_intro(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Part 2  -  Synthetic Bangladesh microgrid dataset",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)
    _add_text(s, "Why synthetic?",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
              size=18, bold=True, color=TEAL)
    _add_bullets(s, [
        "The CSE side needed a dataset with a time axis and richer features "
        "(THD, frequency, storage state, weather, load)",
        "Real datasets with all that in one place are hard to get",
        "So a physics-informed synthetic dataset was generated to test the pipeline",
    ], Inches(0.8), Inches(2.0), Inches(12), Inches(1.8),
       size=15, color=DARK, spacing_after=8)

    warn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.7), Inches(4.2),
                              Inches(11.9), Inches(1.4))
    warn.fill.solid(); warn.fill.fore_color.rgb = CORAL
    warn.line.color.rgb = CORAL
    tf = warn.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = "IMPORTANT  -  synthetic data is for pipeline testing only"
    run.font.size = Pt(15); run.font.bold = True; run.font.color.rgb = WHITE
    for line in [
        ("Every row has  source = SYNTHETIC_GENERATOR_v1  watermark"),
        ("Before publication, real Simulink output from the IEEE 14-bus "
         "model will replace this data"),
    ]:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r = p2.add_run(); r.text = "  •  " + line
        r.font.size = Pt(12); r.font.color.rgb = WHITE

    _add_text(s,
              "Dataset:  50,000 rows  x  50 columns  |  "
              "10 operating scenarios  |  34-day span at 1-sample-per-minute",
              Inches(0.7), Inches(5.9), Inches(12), Inches(0.6),
              size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    _add_footer(s, "Synthetic dataset", 10)


def slide_synth_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Synthetic dataset  -  50 columns of physics",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    left_col = [
        "Time (3):    timestamp, hour_of_day, day_of_year",
        "Environment (4):    irradiance, temp, humidity, wind_speed",
        "Voltage (4):    V_rms_a/b/c, V_unbalance_pct",
        "Current (4):    I_rms_a/b/c, I_neutral",
        "Power (4):    P_active, Q_reactive, S_apparent, PF",
        "Frequency (3):    freq_Hz, freq_dev, RoCoF",
    ]
    right_col = [
        "Harmonics (6):    V-THD, I-THD, 5th, 7th, 11th, 13th",
        "Renewables (3):    PV_kW, wind_kW, RE_penetration",
        "Storage (6):    SOC, SOH, V, I, T, P",
        "Load (2):    load_kW, nonlinear_load_frac",
        "Mode flags (5):    grid_connected, APF, ESS, AI, fault",
        "Labels + economics + watermark (7)",
    ]
    _add_bullets(s, left_col, Inches(0.7), Inches(1.6), Inches(6.2),
                 Inches(4.8), size=12, color=DARK, spacing_after=8)
    _add_bullets(s, right_col, Inches(7.0), Inches(1.6), Inches(6.2),
                 Inches(4.8), size=12, color=DARK, spacing_after=8)

    _add_text(s,
              "Physical relationships built in (not independent noise):",
              Inches(0.7), Inches(6.0), Inches(12), Inches(0.4),
              size=13, bold=True, color=TEAL)
    _add_text(s,
              "PV = f(irradiance, NOCT temp)  |  wind = f(speed^3)  |  "
              "V-THD grows with RE-penetration & nonlinear-load  |  "
              "battery T rises with I-THD^2  |  IEEE 519 verdict from V-THD + I-THD",
              Inches(0.7), Inches(6.4), Inches(12), Inches(0.9),
              size=11, color=GRAY)
    _add_footer(s, "Synthetic schema", 11)


def slide_compliance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "IEEE 519 compliance classifier  (3-class)",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_text(s,
              "Predict PASS / MARGINAL / FAIL from OPERATIONAL SIGNALS ALONE "
              "(V-THD and harmonics are dropped to avoid trivial prediction)",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.9),
              size=14, color=DARK)

    _add_table(s,
               headers=["Model", "Accuracy", "Macro-F1", "FAIL-F1"],
               rows=[
                   ["XGBoost (winner)", "99.82 %", "99.55 %", "99.28 %"],
                   ["Random Forest",     "99.81 %", "99.41 %", "98.86 %"],
                   ["HistGB",            "99.79 %", "99.36 %", "98.79 %"],
                   ["Logistic Regression","99.67 %", "98.90 %", "97.81 %"],
                   ["MLP",               "99.18 %", "97.20 %", "94.41 %"],
               ],
               left=Inches(0.8), top=Inches(2.5),
               width=Inches(6.2), height=Inches(2.5),
               header_color=TEAL, body_size=12, header_size=13)

    _add_image(s, FIG_DIR / "compliance" / "01_confusion_matrix.png",
               Inches(7.3), Inches(2.3), width=Inches(5.5))

    honest = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.7), Inches(5.4),
                                Inches(11.9), Inches(1.4))
    honest.fill.solid(); honest.fill.fore_color.rgb = CORAL
    honest.line.color.rgb = CORAL
    tf = honest.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Honest note  -  these numbers are too clean to be real"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r = p2.add_run()
    r.text = ("The synthetic physics is deterministic.  Once real Simulink data "
              "replaces the synthetic set, accuracies will drop into the "
              "85-92 % range - a much more publishable story.")
    r.font.size = Pt(11); r.font.color.rgb = WHITE
    _add_footer(s, "Compliance classifier", 12)


def slide_forecasting(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Multi-horizon V-THD forecasting",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_text(s,
              "Predict V-THD at 5, 15, and 30 minutes ahead using a "
              "30-minute lookback of 27 operational signals.",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.7),
              size=13, color=DARK)

    _add_table(s,
               headers=["Model", "RMSE (5 min)", "MAE", "R^2", "Early-warn F1"],
               rows=[
                   ["Persistence (naive)", "1.50", "1.17", "-0.67", "0.088"],
                   ["MLP",                 "1.07", "0.84", "+0.15", "0.000"],
                   ["LSTM (winner)",       "1.06", "0.83", "+0.16", "0.000"],
                   ["GRU",                 "1.07", "0.83", "+0.16", "0.000"],
               ],
               left=Inches(0.7), top=Inches(2.4),
               width=Inches(6.5), height=Inches(2.0),
               header_color=TEAL, body_size=12, header_size=13)

    _add_image(s, FIG_DIR / "forecasting" / "02_RMSE_bars.png",
               Inches(7.5), Inches(2.2), width=Inches(5.4))

    _add_text(s, "Two findings",
              Inches(0.7), Inches(4.8), Inches(12), Inches(0.5),
              size=15, bold=True, color=TEAL)
    _add_bullets(s, [
        "Good: LSTM / GRU / MLP all beat persistence by ~30 % on RMSE  "
        "-  the pipeline extracts real signal",
        "Honest: early-warning F1 = 0 because breaches are only 2 %; MSE "
        "loss pushes predictions toward the mean.  Next: reframe as binary "
        "breach classification with class-weighted loss.",
    ], Inches(0.8), Inches(5.2), Inches(12), Inches(1.8),
       size=12, color=DARK, spacing_after=8)
    _add_footer(s, "Forecasting", 13)


def slide_eee(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "EEE-side physics simulation",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_text(s, "Bangladesh microgrid digital twin",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
              size=18, bold=True, color=TEAL)
    _add_bullets(s, [
        "3-phase 50 Hz / 230 V source  +  6-pulse rectifier nonlinear load",
        "SRF (synchronous reference frame) Active Power Filter",
        "FFT-based THD analyzer over a 5-cycle window",
        "Battery-degradation twin  ->  Bangladesh Taka annualised cost",
        "Runs in both MATLAB Simulink AND Python (identical math)",
    ], Inches(0.8), Inches(2.0), Inches(12), Inches(2.5),
       size=14, color=DARK, spacing_after=8)

    _add_image(s, FIG_DIR / "eee_sim" / "02_thd_timeseries.png",
               Inches(0.7), Inches(4.5), width=Inches(6.4))
    _add_image(s, FIG_DIR / "eee_sim" / "04_spectrum_before_after.png",
               Inches(7.2), Inches(4.5), width=Inches(5.7))

    _add_text(s,
              "Headline result:  THD_i  27.13 %  ->  2.07 %  when the APF turns on",
              Inches(0.7), Inches(6.9), Inches(12), Inches(0.4),
              size=13, bold=True, color=CORAL, align=PP_ALIGN.CENTER)
    _add_footer(s, "EEE simulation", 14)


def slide_bridge(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _add_text(s, "How the two halves meet",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=32, bold=True, color=WHITE)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CREAM, weight=2.5)

    _add_text(s,
              "The EEE compliance table has a scenario called "
              '"AI-Optimized Control" that achieves the lowest THD '
              "at the highest renewable penetration.",
              Inches(0.7), Inches(1.7), Inches(12), Inches(1.2),
              size=17, color=LIGHT)

    _add_text(s,
              "That AI is the CSE side's trained ML model.",
              Inches(0.7), Inches(3.0), Inches(12), Inches(0.6),
              size=22, bold=True, color=CREAM, align=PP_ALIGN.CENTER)

    _add_bullets(s, [
        "Without the bridge  =  two separate boring papers "
        "(ML paper + power-systems paper)",
        "With the bridge  =  one interdisciplinary paper where the ML "
        "controls the physical microgrid",
        "This is what makes the work suitable for a Q1 venue rather "
        "than two lower-tier ones",
    ], Inches(0.8), Inches(4.0), Inches(12), Inches(2.5),
       size=14, color=LIGHT, spacing_after=12)
    _add_footer(s, "Interdisciplinary bridge", 15)


def slide_honesty(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "Honesty notes  -  what we do NOT claim",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_text(s,
              "Every claim is bounded up-front so no reviewer can accuse "
              "the paper of overreach.",
              Inches(0.7), Inches(1.5), Inches(12), Inches(0.7),
              size=15, color=GRAY)

    _add_bullets(s, [
        "UCI Grid is static configurations, not time series - we do NOT "
        "claim 'early warning' prediction on it",
        "Results are not claimed to transfer directly to real PMU / SCADA data",
        "tau-robustness margin is a property of the classifier's decision, "
        "not a guarantee about the physical grid",
        "Synthetic dataset carries a source watermark and cannot be used as "
        "evidence in publication tables",
        "Point-forecast MSE + 2 % breach rate has known early-warning "
        "limitations - documented, not hidden",
    ], Inches(0.8), Inches(2.4), Inches(12), Inches(4.0),
       size=14, color=DARK, spacing_after=10)

    _add_text(s,
              "Faisal Sir's Bangla research-integrity announcement explicitly "
              "listed Fabricated Data - so we watermark, disclose, and plan "
              "to replace synthetic with real Simulink output.",
              Inches(0.7), Inches(6.4), Inches(12), Inches(0.6),
              size=12, color=CORAL, align=PP_ALIGN.CENTER)
    _add_footer(s, "Honesty notes", 16)


def slide_gaps(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _add_text(s, "What still needs to happen",
              Inches(0.7), Inches(0.5), Inches(12), Inches(0.7),
              size=28, bold=True, color=DARK)
    _add_line(s, Inches(0.7), Inches(1.3), Inches(12.6), Inches(1.3),
              color=CORAL, weight=2.5)

    _add_text(s, "CSE-side next steps",
              Inches(0.7), Inches(1.6), Inches(12), Inches(0.5),
              size=18, bold=True, color=TEAL)
    _add_bullets(s, [
        "Class-weighted or binary reformulation of the forecaster",
        "Multi-seed + bootstrap CI on all forecasting numbers",
        "Reproducibility audit for the synthetic-data side",
        "Draft the LaTeX manuscript from PAPER_OUTLINE.md",
    ], Inches(0.8), Inches(2.1), Inches(12), Inches(1.8),
       size=13, color=DARK, spacing_after=6)

    _add_text(s, "EEE-side next steps",
              Inches(0.7), Inches(4.1), Inches(12), Inches(0.5),
              size=18, bold=True, color=CORAL)
    _add_bullets(s, [
        "Finish the IEEE 14-bus parametric-study Simulink runs",
        "Formal Schafer 2016 reproduction on the 4-bus model",
        "Complete the AI Forecasting Performance and EMS Optimisation sheets",
        "Send real .slx output files to the CSE side to replace synthetic data",
    ], Inches(0.8), Inches(4.6), Inches(12), Inches(2.0),
       size=13, color=DARK, spacing_after=6)
    _add_footer(s, "Gaps and TODO", 17)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(2.6),
                             Inches(0.15), Inches(2.2))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    _add_text(s, "Thanks  -  questions welcome",
              Inches(1.0), Inches(2.5), Inches(12), Inches(1.2),
              size=44, bold=True, color=WHITE)
    _add_text(s,
              "Repo:  github.com/Fahmidaca/microgrid-dt-1",
              Inches(1.0), Inches(3.7), Inches(12), Inches(0.6),
              size=20, color=TEAL, font="Consolas")
    _add_text(s,
              "Target venues:  IEEE Trans. Smart Grid (Q1, IF~10)  |  "
              "Applied Energy (Q1, IF~11)",
              Inches(1.0), Inches(4.5), Inches(12), Inches(0.6),
              size=15, color=CREAM)
    _add_text(s,
              "Every folder in the repo has a README explaining what it "
              "contains and how to reproduce it.",
              Inches(1.0), Inches(5.4), Inches(12), Inches(0.6),
              size=13, color=LIGHT)


# =========================================================================
#  DRIVER
# =========================================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_two_halves(prs)
    slide_uci_dataset(prs)
    slide_uci_methodology(prs)
    slide_uci_results(prs)
    slide_tau_margin(prs)
    slide_tau_insight(prs)
    slide_synth_intro(prs)
    slide_synth_features(prs)
    slide_compliance(prs)
    slide_forecasting(prs)
    slide_eee(prs)
    slide_bridge(prs)
    slide_honesty(prs)
    slide_gaps(prs)
    slide_thanks(prs)

    out_path = OUT_DIR / "microgrid_paper_progress.pptx"
    prs.save(out_path)
    print(f"[pptx] wrote {out_path} ({out_path.stat().st_size / 1024:.1f} KB)")
    print(f"[pptx] {len(prs.slides)} slides")
    return out_path


if __name__ == "__main__":
    build()
