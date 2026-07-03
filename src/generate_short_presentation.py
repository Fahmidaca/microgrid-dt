"""
Generate microgrid_paper_progress_short.pptx  -  an 8-slide condensed
version of the full 18-slide deck.

Slide flow (each slide combines material from 2-3 slides of the full deck):
    1. Title
    2. Problem + two-halves project layout
    3. Part 1  UCI Grid Stability  (dataset + methodology + results in one)
    4. Novel contribution  tau-robustness margin  (metric + insight)
    5. Part 2  Synthetic dataset  (intro + schema + watermark)
    6. Downstream ML results  (compliance + forecasting side by side)
    7. EEE simulation + the bridge  (THD result + AI-Optimized Control)
    8. Honesty + gaps + thanks

Run:
    python src/generate_short_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"; OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR = ROOT / "figures"

# --- theme (matches the long deck)
DARK  = RGBColor(0x26, 0x46, 0x53)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
CREAM = RGBColor(0xE9, 0xC4, 0x6A)
LIGHT = RGBColor(0xF4, 0xF1, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x60, 0x60, 0x60)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def _paint_bg(slide, color=DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background(); bg.fill.solid()
    bg.fill.fore_color.rgb = color; bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element); spTree.insert(2, bg._element)
    return bg


def _text(slide, txt, left, top, width, height, size=18,
          bold=False, color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def _bullets(slide, items, left, top, width, height, size=13,
             color=LIGHT, spacing=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(spacing)
        r = p.add_run(); r.text = f"•  {b}"
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _line(slide, x1, y1, x2, y2, color=CORAL, weight=2.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln


def _image(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), left, top,
                                        width=width, height=height)
    return None


def _footer(slide, label, page):
    _text(slide, label, Inches(0.4), Inches(7.05), Inches(9),
          Inches(0.35), size=9, color=GRAY)
    _text(slide, str(page), Inches(12.5), Inches(7.05), Inches(0.5),
          Inches(0.35), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def _table(slide, headers, rows, left, top, width, height,
           header_color=TEAL, body_size=11, header_size=11):
    n_rows = len(rows) + 1; n_cols = len(headers)
    shp = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shp.table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        c.text_frame.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = (WHITE if i % 2 == 1
                                     else RGBColor(0xF0, 0xF0, 0xF0))
            c.text_frame.text = ""
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(body_size); r.font.color.rgb = DARK
            r.font.name = "Calibri"
    return shp


def _rounded_box(slide, left, top, width, height, fill, text_lines,
                 first_bold=True, first_size=15, body_size=11,
                 text_color=WHITE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.10)
    for i, line in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        if i == 0 and first_bold:
            r.font.bold = True; r.font.size = Pt(first_size)
        else:
            r.font.size = Pt(body_size)
        r.font.color.rgb = text_color; r.font.name = "Calibri"
    return box


# ============ SLIDES =====================================================

def slide_1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    # Coral accent bar spans the whole text stack
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(1.7), Inches(0.15), Inches(3.9))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL

    _text(s, "Microgrid Digital Twin",
          Inches(1.0), Inches(1.7), Inches(11.5), Inches(1.2),
          size=54, bold=True, color=WHITE)
    _text(s, "Stability, Power Quality, and Forecasting",
          Inches(1.0), Inches(2.95), Inches(11.5), Inches(0.9),
          size=34, color=TEAL)
    _text(s, "Condensed progress report  -  CSE + EEE sides",
          Inches(1.0), Inches(3.95), Inches(11.5), Inches(0.7),
          size=24, bold=True, color=CREAM)
    _text(s,
          "Target venues:  IEEE Trans. Smart Grid (Q1, IF~10)  |  "
          "Applied Energy (Q1, IF~11)",
          Inches(1.0), Inches(4.9), Inches(11.5), Inches(0.6),
          size=19, color=LIGHT)
    _text(s, "github.com/Fahmidaca/microgrid-dt-1",
          Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.5),
          size=17, color=CREAM, font="Consolas")


def slide_2_problem_and_layout(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "The problem  +  project layout",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    # LEFT: the 5 questions
    _text(s, "5 questions this project answers",
          Inches(0.7), Inches(1.35), Inches(6), Inches(0.5),
          size=16, bold=True, color=TEAL)
    _bullets(s, [
        "Is the microgrid stable right now?",
        "How much can we trust the model that says yes?",
        "Will THD breach IEEE 519's 5 % limit in the next 5-30 minutes?",
        "Do renewable-heavy scenarios pass IEEE 519 without help?",
        "Can an AI controller beat a fixed control law?",
    ], Inches(0.7), Inches(1.85), Inches(6), Inches(3.5),
       size=13, color=DARK, spacing=8)

    # RIGHT: two-halves diagram
    _text(s, "Two halves that meet in the middle",
          Inches(7.0), Inches(1.35), Inches(6), Inches(0.5),
          size=16, bold=True, color=CORAL)

    # Left half (CSE)
    _rounded_box(s, Inches(7.0), Inches(1.85), Inches(2.85), Inches(3.5),
                 fill=TEAL, text_lines=[
                     "CSE / ML SIDE",
                     "", "UCI Grid Stability",
                     "5 models x 5 seeds",
                     "MLP  96.05 %",
                     "", "tau-robustness margin",
                     "(new metric)",
                     "", "Synthetic 50k dataset",
                     "Compliance + forecasting",
                 ], first_size=13, body_size=10)

    # Right half (EEE)
    _rounded_box(s, Inches(10.0), Inches(1.85), Inches(2.85), Inches(3.5),
                 fill=CORAL, text_lines=[
                     "EEE / PHYSICS SIDE",
                     "", "50 Hz / 230 V microgrid",
                     "6-pulse rectifier",
                     "", "SRF Active Power Filter",
                     "FFT THD analyzer",
                     "Battery ageing + BDT cost",
                     "", "IEEE 14-bus study",
                     "10 scenarios",
                 ], first_size=13, body_size=10)

    # Meeting point across bottom
    meet = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.7), Inches(5.6),
                              Inches(12.1), Inches(1.1))
    meet.fill.solid(); meet.fill.fore_color.rgb = DARK
    meet.line.color.rgb = DARK
    tf = meet.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "MEETING POINT"
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = CREAM
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = ('EEE compliance table has an "AI-Optimized Control" scenario. '
              'That AI is the CSE side\'s trained ML model.')
    r.font.size = Pt(12); r.font.color.rgb = WHITE
    _footer(s, "Problem + layout", 2)


def slide_3_uci(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Part 1  -  UCI Grid Stability pipeline",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    _text(s, "Dataset",
          Inches(0.7), Inches(1.30), Inches(6), Inches(0.4),
          size=15, bold=True, color=TEAL)
    _bullets(s, [
        "Schafer 2016  -  10,000 snapshots of a 4-bus decentralised grid",
        "12 features (tau1..4, p1..4, g1..4)  -  binary stable/unstable label",
        "Peer-reviewed, still cited in 2024-25 (benchmarks don't expire)",
    ], Inches(0.7), Inches(1.75), Inches(6), Inches(1.7),
       size=11, color=DARK, spacing=4)

    _text(s, "Methodology",
          Inches(0.7), Inches(3.65), Inches(6), Inches(0.4),
          size=15, bold=True, color=TEAL)
    _bullets(s, [
        "5 classifiers (LogReg, RF, HistGB, XGBoost, MLP)  x  5 seeds",
        "80/20 stratified split, deterministic seeds",
        "Paired McNemar + Holm-Bonferroni + bootstrap 95 % CI",
        "All 10 pairwise differences significant at p < 0.05",
    ], Inches(0.7), Inches(4.1), Inches(6), Inches(2.5),
       size=11, color=DARK, spacing=4)

    _text(s, "Results  (mean over 5 seeds)",
          Inches(7.2), Inches(1.30), Inches(6), Inches(0.4),
          size=15, bold=True, color=TEAL)
    _table(s,
           headers=["Model", "Accuracy", "95 % CI", "F1"],
           rows=[
               ["MLP  (winner)",  "96.05 %", "[95.65, 96.43]", "95.72 %"],
               ["HistGB",         "94.72 %", "[94.28, 95.13]", "94.24 %"],
               ["XGBoost",        "94.39 %", "[93.92, 94.83]", "93.87 %"],
               ["Random Forest",  "92.16 %", "[91.64, 92.70]", "91.38 %"],
               ["LogReg",         "81.50 %", "[80.72, 82.29]", "79.58 %"],
           ],
           left=Inches(7.2), top=Inches(1.75),
           width=Inches(5.6), height=Inches(2.5),
           body_size=11, header_size=12)

    _text(s,
          "MLP wins on accuracy - but that is only half the story.  "
          "The other half (tau-robustness margin) is on the next slide.",
          Inches(7.2), Inches(4.5), Inches(5.6), Inches(1.5),
          size=12, color=CORAL, align=PP_ALIGN.LEFT)
    _footer(s, "UCI pipeline", 3)


def slide_4_tau(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _text(s, "Novel contribution  -  tau-robustness margin",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=WHITE)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CREAM, weight=2.5)

    _text(s, "What we defined",
          Inches(0.7), Inches(1.35), Inches(6.2), Inches(0.5),
          size=16, bold=True, color=CREAM)
    _bullets(s, [
        "For each correctly-classified test point, binary-search the "
        "smallest reaction-time perturbation that flips the prediction",
        "16 sign patterns on tau1..tau4, clipped to [0.5, 10] s (L_inf norm)",
        "Result: per-sample safety buffer  -  no prior paper reports it",
    ], Inches(0.7), Inches(1.9), Inches(6.2), Inches(2.5),
       size=12, color=LIGHT, spacing=6)

    _text(s, "Key insight  -  accuracy != robustness",
          Inches(0.7), Inches(4.4), Inches(6.2), Inches(0.5),
          size=16, bold=True, color=CORAL)
    _bullets(s, [
        "MLP:  highest accuracy (96 %)  BUT smallest tau-margin",
        "LogReg:  lowest accuracy (81.5 %)  BUT largest margin  "
        "(stable 1.33 s / unstable 2.12 s)",
        "Prior accuracy-only reports hide this trade-off",
        "'More accurate' != 'more deployable' - the publishable insight",
    ], Inches(0.7), Inches(4.95), Inches(6.2), Inches(2.2),
       size=12, color=LIGHT, spacing=5)

    _image(s, FIG_DIR / "robustness" / "02_margin_per_class.png",
           Inches(7.5), Inches(1.5), width=Inches(5.4))
    _footer(s, "Novel contribution", 4)


def slide_5_synthetic(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Part 2  -  synthetic Bangladesh microgrid dataset",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    _text(s, "50,000 rows  x  50 columns  |  34 days  |  1 sample per minute",
          Inches(0.7), Inches(1.3), Inches(12), Inches(0.4),
          size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # left: schema summary
    _text(s, "50 columns in 12 categories",
          Inches(0.7), Inches(1.9), Inches(6.2), Inches(0.4),
          size=13, bold=True, color=DARK)
    _bullets(s, [
        "Time (3), Environment (4), Voltage (4), Current (4)",
        "Power (4), Frequency (3), Harmonics (6), Renewables (3)",
        "Storage (6), Load (2), Mode flags (5), Labels + cost + watermark (6)",
        "Physics-informed:  PV = f(irradiance, NOCT temp),  "
        "wind = f(speed^3)",
        "THD grows with RE penetration + nonlinear load,  falls with "
        "ESS + AI control",
    ], Inches(0.7), Inches(2.3), Inches(6.2), Inches(2.7),
       size=11, color=DARK, spacing=4)

    # right: watermark warning
    warn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.2), Inches(1.9),
                              Inches(5.6), Inches(3.0))
    warn.fill.solid(); warn.fill.fore_color.rgb = CORAL
    warn.line.color.rgb = CORAL
    tf = warn.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "IMPORTANT  -  pipeline testing only"
    r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = WHITE
    for line in [
        "",
        "Every row watermarked with",
        '     source = SYNTHETIC_GENERATOR_v1',
        "",
        "Before publication, real Simulink output from",
        "the EEE side's IEEE 14-bus model will replace it.",
        "",
        "(Faisal Sir: fabricated data is a top",
        " research-integrity violation.)",
    ]:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.LEFT
        r = p2.add_run(); r.text = line
        r.font.size = Pt(11); r.font.color.rgb = WHITE

    # Bottom plot strip - constrained by HEIGHT so they don't overflow.
    # Available vertical space: 5.05 -> 6.95 = 1.90 in.  Use 1.55 in tall,
    # centre horizontally on their halves.
    _image(s, FIG_DIR / "synthetic" / "03_daily_profile.png",
           Inches(0.9), Inches(5.10), height=Inches(1.55))
    _image(s, FIG_DIR / "synthetic" / "04_correlations.png",
           Inches(9.5), Inches(5.10), height=Inches(1.55))
    _footer(s, "Synthetic dataset", 5)


def slide_6_downstream(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Downstream ML on the synthetic dataset",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    # LEFT: compliance classifier
    _text(s, "IEEE 519 compliance  (PASS / MARGINAL / FAIL)",
          Inches(0.5), Inches(1.30), Inches(6.3), Inches(0.4),
          size=14, bold=True, color=TEAL)
    _text(s,
          "43 features  |  V-THD & harmonics dropped as leakage  "
          "|  5 models x 5 seeds",
          Inches(0.5), Inches(1.75), Inches(6.3), Inches(0.4),
          size=10, color=GRAY)
    _table(s,
           headers=["Model", "Accuracy", "Macro-F1"],
           rows=[
               ["XGBoost (winner)", "99.82 %", "99.55 %"],
               ["Random Forest",    "99.81 %", "99.41 %"],
               ["HistGB",           "99.79 %", "99.36 %"],
               ["LogReg",           "99.67 %", "98.90 %"],
               ["MLP",              "99.18 %", "97.20 %"],
           ],
           left=Inches(0.5), top=Inches(2.15),
           width=Inches(6.3), height=Inches(2.3),
           body_size=11, header_size=11)

    _rounded_box(s, Inches(0.5), Inches(4.7),
                 Inches(6.3), Inches(1.9),
                 fill=CORAL, text_lines=[
                     "Honest note",
                     "",
                     "Numbers are too clean - synthetic physics is",
                     "deterministic, so features like RE penetration",
                     "essentially determine the label.",
                     "",
                     "Real Simulink data will drop accuracies into",
                     "the 85-92 % range  -  a better paper story.",
                 ], first_size=13, body_size=10)

    # RIGHT: forecasting
    _text(s, "Multi-horizon V-THD forecasting",
          Inches(7.0), Inches(1.30), Inches(6.0), Inches(0.4),
          size=14, bold=True, color=TEAL)
    _text(s,
          "30-min lookback  ->  predict at 5 / 15 / 30 min ahead  "
          "|  chronological split",
          Inches(7.0), Inches(1.75), Inches(6.0), Inches(0.4),
          size=10, color=GRAY)
    _table(s,
           headers=["Model (5 min)", "RMSE", "R^2", "EW-F1"],
           rows=[
               ["Persistence",     "1.50", "-0.67", "0.088"],
               ["MLP",             "1.07", "+0.15", "0.000"],
               ["LSTM (winner)",   "1.06", "+0.16", "0.000"],
               ["GRU",             "1.07", "+0.16", "0.000"],
           ],
           left=Inches(7.0), top=Inches(2.15),
           width=Inches(6.0), height=Inches(2.0),
           body_size=11, header_size=11)

    _rounded_box(s, Inches(7.0), Inches(4.7),
                 Inches(6.0), Inches(1.9),
                 fill=CORAL, text_lines=[
                     "Honest note",
                     "",
                     "Deep models beat persistence by 30 % on RMSE",
                     "-  the pipeline extracts real signal.",
                     "",
                     "BUT early-warning F1 = 0 (breaches are only 2 %,",
                     "MSE hugs the mean).  Next step:  binary breach",
                     "classification with class-weighted loss.",
                 ], first_size=13, body_size=10)

    _footer(s, "Downstream ML", 6)


def slide_7_eee_bridge(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "EEE physics simulation  +  the bridge",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    # LEFT: EEE simulation
    _text(s, "Bangladesh PV microgrid digital twin",
          Inches(0.7), Inches(1.3), Inches(6.2), Inches(0.5),
          size=15, bold=True, color=TEAL)
    _bullets(s, [
        "3-phase 50 Hz / 230 V + 6-pulse rectifier",
        "SRF Active Power Filter + FFT THD analyzer",
        "Battery ageing + BDT cost twin",
        "Runs in MATLAB Simulink AND Python (same math)",
    ], Inches(0.7), Inches(1.85), Inches(6.2), Inches(2.0),
       size=11, color=DARK, spacing=4)
    _image(s, FIG_DIR / "eee_sim" / "02_thd_timeseries.png",
           Inches(0.7), Inches(4.0), width=Inches(6.2))
    _text(s, "Result:  THD_i  27.1 %  ->  2.1 %  when APF turns on",
          Inches(0.7), Inches(6.35), Inches(6.2), Inches(0.4),
          size=13, bold=True, color=CORAL, align=PP_ALIGN.CENTER)

    # RIGHT: the bridge
    _text(s, "How the two halves meet",
          Inches(7.2), Inches(1.3), Inches(5.6), Inches(0.5),
          size=15, bold=True, color=CORAL)
    _text(s,
          "The EEE compliance table has a scenario called "
          "'AI-Optimized Control' that achieves the best THD.",
          Inches(7.2), Inches(1.85), Inches(5.6), Inches(1.0),
          size=12, color=DARK)
    _rounded_box(s, Inches(7.2), Inches(3.0),
                 Inches(5.6), Inches(1.2),
                 fill=DARK, text_lines=[
                     "That AI is the CSE side's",
                     "trained ML model.",
                 ], first_size=15, body_size=15, text_color=CREAM)

    _text(s, "Why it matters",
          Inches(7.2), Inches(4.4), Inches(5.6), Inches(0.4),
          size=13, bold=True, color=TEAL)
    _bullets(s, [
        "Without the bridge = two boring papers",
        "With the bridge = one interdisciplinary paper",
        "The ML controls the physical microgrid,  "
        "Simulink physics validates it beats fixed control laws",
        "Unlocks Q1 venues instead of two mid-tier ones",
    ], Inches(7.2), Inches(4.85), Inches(5.6), Inches(2.0),
       size=10, color=DARK, spacing=4)

    _footer(s, "EEE + bridge", 7)


def slide_8_wrap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _text(s, "Honesty  +  gaps  +  thanks",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=WHITE)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CREAM, weight=2.5)

    # Honesty column
    _text(s, "Honesty notes  -  what we do NOT claim",
          Inches(0.7), Inches(1.3), Inches(6.0), Inches(0.5),
          size=15, bold=True, color=CREAM)
    _bullets(s, [
        "UCI is static  -  no 'early warning' claim on it",
        "No claim of transfer to real PMU / SCADA",
        "Tau-margin is a decision property, not a grid guarantee",
        "Synthetic dataset carries a watermark; not for publication",
        "Forecaster's zero early-warning F1 is a real limitation, not hidden",
    ], Inches(0.7), Inches(1.85), Inches(6.0), Inches(3.0),
       size=11, color=LIGHT, spacing=6)

    # Gaps column
    _text(s, "What still needs to happen",
          Inches(7.0), Inches(1.3), Inches(6.0), Inches(0.5),
          size=15, bold=True, color=CORAL)
    _text(s, "CSE side",
          Inches(7.0), Inches(1.85), Inches(6.0), Inches(0.35),
          size=12, bold=True, color=TEAL)
    _bullets(s, [
        "Binary breach classifier with class-weighted loss",
        "Multi-seed + bootstrap CI on forecasting",
        "Reproducibility audit on synthetic side",
        "Draft the LaTeX manuscript",
    ], Inches(7.0), Inches(2.25), Inches(6.0), Inches(1.7),
       size=10, color=LIGHT, spacing=3)

    _text(s, "EEE side",
          Inches(7.0), Inches(3.95), Inches(6.0), Inches(0.35),
          size=12, bold=True, color=TEAL)
    _bullets(s, [
        "Finish IEEE 14-bus parametric-study runs",
        "Formal Schafer 2016 reproduction on 4-bus model",
        "Send real .slx output to replace synthetic data",
    ], Inches(7.0), Inches(4.35), Inches(6.0), Inches(1.4),
       size=10, color=LIGHT, spacing=3)

    # Thanks strip across bottom
    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.7), Inches(5.7),
                                Inches(12.1), Inches(1.15))
    strip.fill.solid(); strip.fill.fore_color.rgb = CORAL
    strip.line.color.rgb = CORAL
    tf = strip.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Thanks  -  questions welcome"
    r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run()
    r.text = "github.com/Fahmidaca/microgrid-dt-1"
    r.font.size = Pt(13); r.font.color.rgb = CREAM; r.font.name = "Consolas"

    _footer(s, "Wrap-up", 8)


# ============ DRIVER =====================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    slide_1_title(prs)
    slide_2_problem_and_layout(prs)
    slide_3_uci(prs)
    slide_4_tau(prs)
    slide_5_synthetic(prs)
    slide_6_downstream(prs)
    slide_7_eee_bridge(prs)
    slide_8_wrap(prs)
    out = OUT_DIR / "microgrid_paper_progress_short.pptx"
    prs.save(out)
    print(f"[pptx-short] wrote {out} ({out.stat().st_size/1024:.1f} KB, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
