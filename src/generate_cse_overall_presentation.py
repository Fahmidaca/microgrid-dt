"""
Generate cse_overall_progress.pptx  -  8-slide CSE-only progress deck
covering everything the CSE side has built, old and new.

Slide flow:
    1. Title
    2. CSE scope  -  3 datasets, what's built on each
    3. UCI Grid Stability  +  tau-robustness margin
    4. Synthetic dataset  +  IEEE 519 compliance  +  THD forecasting
    5. Anomaly detection  -  the two-dataset comparison (today's centerpiece)
    6. Disturbance prediction  +  explainability
    7. Dashboard  +  data provenance resolution
    8. Status checklist  +  thanks

Uses a mix of chart types on purpose: embedded PNG figures from the
actual pipeline runs, a native (editable) PowerPoint bar chart, and
data tables - not just one format repeated eight times.

Run:
    python src/generate_cse_overall_presentation.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "presentation"; OUT_DIR.mkdir(parents=True, exist_ok=True)
FIG_DIR = ROOT / "figures"

# --- theme (matches the other decks in this repo)
DARK  = RGBColor(0x26, 0x46, 0x53)
TEAL  = RGBColor(0x2A, 0x9D, 0x8F)
CORAL = RGBColor(0xE7, 0x6F, 0x51)
CREAM = RGBColor(0xE9, 0xC4, 0x6A)
LIGHT = RGBColor(0xF4, 0xF1, 0xDE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0x60, 0x60, 0x60)

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)


def _paint_bg(slide, color=DARK):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background(); bg.fill.solid()
    bg.fill.fore_color.rgb = color; bg.shadow.inherit = False
    spTree = slide.shapes._spTree
    spTree.remove(bg._element); spTree.insert(2, bg._element)
    return bg


def _text(slide, txt, left, top, width, height, size=18,
          bold=False, color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = txt
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def _bullets(slide, items, left, top, width, height, size=13,
             color=LIGHT, spacing=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, b in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT; p.space_after = Pt(spacing)
        r = p.add_run(); r.text = f"•  {b}"
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def _line(slide, x1, y1, x2, y2, color=CORAL, weight=2.0):
    ln = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln.line.color.rgb = color; ln.line.width = Pt(weight)
    return ln


def _image(slide, path, left, top, width=None, height=None):
    if Path(path).exists():
        return slide.shapes.add_picture(str(path), left, top,
                                        width=width, height=height)
    return None


def _footer(slide, label, page):
    _text(slide, label, Inches(0.4), Inches(7.05), Inches(9),
          Inches(0.35), size=9, color=GRAY)
    _text(slide, str(page), Inches(12.5), Inches(7.05), Inches(0.5),
          Inches(0.35), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def _table(slide, headers, rows, left, top, width, height,
           header_color=TEAL, body_size=11, header_size=11):
    n_rows = len(rows) + 1; n_cols = len(headers)
    shp = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shp.table
    for j, h in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = header_color
        c.text_frame.text = ""
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = h
        r.font.size = Pt(header_size); r.font.bold = True
        r.font.color.rgb = WHITE; r.font.name = "Calibri"
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = (WHITE if i % 2 == 1
                                     else RGBColor(0xF0, 0xF0, 0xF0))
            c.text_frame.text = ""
            p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(v)
            r.font.size = Pt(body_size); r.font.color.rgb = DARK
            r.font.name = "Calibri"
    return shp


def _rounded_box(slide, left, top, width, height, fill, text_lines,
                 first_bold=True, first_size=15, body_size=11,
                 text_color=WHITE):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.10)
    for i, line in enumerate(text_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        if i == 0 and first_bold:
            r.font.bold = True; r.font.size = Pt(first_size)
        else:
            r.font.size = Pt(body_size)
        r.font.color.rgb = text_color; r.font.name = "Calibri"
    return box


def _native_bar_chart(slide, categories, series_name, values, left, top,
                       width, height, color=TEAL):
    """A real, editable PowerPoint chart object - not an image."""
    data = CategoryChartData()
    data.categories = categories
    data.add_series(series_name, values)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, data)
    chart = gframe.chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dls = plot.data_labels
    dls.number_format = '0.0"%"'; dls.number_format_is_linked = False
    dls.font.size = Pt(11); dls.font.color.rgb = DARK
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.visible = False
    return gframe


# ============ SLIDES =====================================================

def slide_1_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.7), Inches(1.6), Inches(0.15), Inches(4.1))
    bar.line.fill.background(); bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL

    _text(s, "CSE-Side Progress",
          Inches(1.0), Inches(1.6), Inches(11.5), Inches(1.2),
          size=54, bold=True, color=WHITE)
    _text(s, "Microgrid Digital Twin — Stability, Power Quality, and Forecasting",
          Inches(1.0), Inches(2.85), Inches(11.5), Inches(0.7),
          size=26, color=TEAL)
    _text(s,
          "Stability  •  Compliance  •  Forecasting  •  "
          "Cyber-Resilience  •  Explainability  •  Dashboard",
          Inches(1.0), Inches(3.65), Inches(11.5), Inches(0.7),
          size=20, bold=True, color=CREAM)
    _text(s,
          "5 datasets-to-model pipelines, all owned end-to-end by the CSE side",
          Inches(1.0), Inches(4.55), Inches(11.5), Inches(0.5),
          size=15, color=LIGHT)
    _text(s, "github.com/Fahmidaca/microgrid-dt",
          Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.5),
          size=17, color=CREAM, font="Consolas")


def slide_2_scope(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "CSE scope — 3 datasets, 5 pipelines",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    boxes = [
        (TEAL, ["UCI Grid Stability", "(public benchmark)", "",
                "10,000 rows, 12 features", "5 classifiers x 5 seeds",
                "+ tau-robustness margin"]),
        (CORAL, ["Synthetic Bangladesh", "microgrid dataset", "",
                 "50,000 rows x 50 cols", "IEEE 519 compliance classifier",
                 "Multi-horizon THD forecasting"]),
        (DARK, ["Power-quality disturbance", "dataset (Jamalpur, field data)", "",
                "5,000 rows x 16 cols", "Anomaly detection + prediction",
                "SHAP explainability + dashboard"]),
    ]
    x = Inches(0.7)
    for fill, lines in boxes:
        _rounded_box(s, x, Inches(1.5), Inches(3.95), Inches(3.9),
                     fill=fill, text_lines=lines, first_size=15, body_size=11)
        x += Inches(4.15)

    _rounded_box(s, Inches(0.7), Inches(5.6), Inches(12.1), Inches(1.1),
                 fill=CREAM, text_lines=[
                     "Two of the three datasets were compared side-by-side on the "
                     "same anomaly-detection method rather than merged — "
                     "incompatible schemas, but a very useful comparison (slide 5).",
                 ], first_bold=False, body_size=13, text_color=DARK)
    _footer(s, "CSE scope", 2)


def slide_3_uci(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "UCI Grid Stability  +  tau-robustness margin",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    _table(s,
           headers=["Model", "Accuracy", "F1"],
           rows=[
               ["MLP (winner)", "96.05%", "95.72%"],
               ["HistGB", "94.72%", "94.24%"],
               ["XGBoost", "94.39%", "93.87%"],
               ["RF", "92.16%", "91.38%"],
               ["LogReg", "81.50%", "79.58%"],
           ],
           left=Inches(0.7), top=Inches(1.4),
           width=Inches(5.6), height=Inches(2.4),
           body_size=12, header_size=12)

    _rounded_box(s, Inches(0.7), Inches(4.0), Inches(5.6), Inches(2.7),
                 fill=DARK, text_lines=[
                     "Novel: tau-robustness margin",
                     "",
                     "Per-sample: smallest reaction-time change",
                     "that flips the model's prediction.",
                     "",
                     "MLP: highest accuracy, SMALLEST margin.",
                     "LogReg: lowest accuracy, LARGEST margin.",
                     "'More accurate' != 'more deployable'.",
                 ], first_size=15, body_size=12, text_color=CREAM)

    _image(s, FIG_DIR / "robustness" / "02_margin_per_class.png",
           Inches(6.6), Inches(1.4), width=Inches(6.3))
    _footer(s, "UCI + robustness", 3)


def slide_4_synthetic_downstream(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Synthetic dataset  +  compliance  +  forecasting",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    _text(s, "IEEE 519 compliance classifier (PASS/MARGINAL/FAIL)",
          Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.4),
          size=13, bold=True, color=TEAL)
    _native_bar_chart(
        s, ["XGBoost", "RF", "HistGB", "LogReg", "MLP"],
        "Accuracy %", [99.82, 99.81, 99.79, 99.67, 99.18],
        Inches(0.7), Inches(1.7), Inches(5.8), Inches(2.5), color=TEAL)

    _text(s, "50,000 rows x 50 cols, physics-informed, watermarked "
             "SYNTHETIC_GENERATOR_v1 — not used as publication evidence",
          Inches(0.7), Inches(4.3), Inches(5.8), Inches(0.6),
          size=10, color=GRAY)

    _rounded_box(s, Inches(0.7), Inches(5.0), Inches(5.8), Inches(1.7),
                 fill=CORAL, text_lines=[
                     "THD forecasting (LSTM/GRU/MLP, 5/15/30 min ahead)",
                     "LSTM beats persistence baseline by ~30% RMSE.",
                     "Honest gap: early-warning F1 = 0 (breaches are 2% "
                     "of rows) — next step is class-weighted loss.",
                 ], first_size=13, body_size=11)

    _image(s, FIG_DIR / "synthetic" / "04_correlations.png",
           Inches(6.7), Inches(1.4), width=Inches(6.2))
    _footer(s, "Synthetic + downstream", 4)


def slide_5_anomaly_comparison(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _text(s, "Anomaly detection — comparing the two datasets",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=WHITE)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CREAM, weight=2.5)

    _image(s, FIG_DIR / "anomaly" / "01_precision_recall_comparison.png",
           Inches(0.6), Inches(1.35), width=Inches(7.0))

    _text(s, "Why compare instead of merge?",
          Inches(7.9), Inches(1.35), Inches(5.0), Inches(0.4),
          size=14, bold=True, color=CREAM)
    _bullets(s, [
        "Incompatible schemas: single-phase 16-col vs. three-phase 50-col",
        "Ran the identical Isolation Forest recipe on both instead",
    ], Inches(7.9), Inches(1.8), Inches(5.0), Inches(1.3),
       size=11, color=LIGHT, spacing=5)

    _text(s, "What it proved",
          Inches(7.9), Inches(3.2), Inches(5.0), Inches(0.4),
          size=14, bold=True, color=CREAM)
    _bullets(s, [
        "Disturbance dataset: 0.70 ROC-AUC — real, learnable fault signal",
        "Synthetic dataset: 0.50 (chance) — CORRECT, its fault_flag is "
        "designed as pure random noise",
        "The null result on one confirms the real result on the "
        "other isn't a fluke",
    ], Inches(7.9), Inches(3.65), Inches(5.0), Inches(2.3),
       size=11, color=LIGHT, spacing=6)

    _footer(s, "Anomaly detection", 5)


def slide_6_prediction_xai(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Disturbance prediction  +  explainability",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    # Both images height-constrained (not width-constrained) so they
    # cannot grow into the honesty box below regardless of aspect ratio.
    _text(s, "4-class model: 99.70% test accuracy",
          Inches(0.7), Inches(1.3), Inches(6.0), Inches(0.4),
          size=13, bold=True, color=TEAL)
    _image(s, FIG_DIR / "disturbance" / "01_confusion_matrix.png",
           Inches(0.9), Inches(1.75), height=Inches(3.55))

    _text(s, "SHAP — why the model predicts what it predicts",
          Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.4),
          size=13, bold=True, color=TEAL)
    _image(s, FIG_DIR / "xai" / "01_shap_summary_bar.png",
           Inches(6.9), Inches(1.9), height=Inches(3.2))

    _rounded_box(s, Inches(0.7), Inches(5.6), Inches(12.1), Inches(1.15),
                 fill=DARK, text_lines=[
                     "voltage_rms_V dominates Voltage_Sag, THD/harmonics dominate "
                     "Harmonic_Distortion — physically sensible, not a black box.",
                 ], first_bold=False, body_size=13, text_color=CREAM)
    _footer(s, "Prediction + XAI", 6)


def slide_7_dashboard_provenance(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, LIGHT)
    _text(s, "Dashboard  +  data provenance — resolved",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=26, bold=True, color=DARK)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CORAL, weight=2.5)

    _text(s, "Interactive dashboard  (streamlit run src/dashboard.py)",
          Inches(0.7), Inches(1.3), Inches(6.0), Inches(0.4),
          size=14, bold=True, color=TEAL)
    _bullets(s, [
        "Pick a row -> live reading (voltage/current/THD/weather)",
        "Alert banners from BOTH the disturbance classifier and the "
        "anomaly detector",
        "Live SHAP explanation for that specific row, not just an "
        "aggregate plot",
        "Cost-impact panel — explicitly labelled as team-calculated, "
        "not measured",
        "Historical trend chart with disturbance rows colour-coded",
    ], Inches(0.7), Inches(1.8), Inches(6.0), Inches(3.5),
       size=12, color=DARK, spacing=7)

    _rounded_box(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(4.0),
                 fill=TEAL, text_lines=[
                     "Data provenance — resolved",
                     "",
                     "Team confirmed: electrical readings are field-",
                     "measured at the Jamalpur powerplant site.",
                     "",
                     "Still true regardless: economic_cost_BDT and",
                     "battery_capacity_loss_pct are team-calculated",
                     "columns, not independent measurements —",
                     "documented column-by-column so the paper",
                     "describes each one correctly.",
                 ], first_size=15, body_size=12)
    _footer(s, "Dashboard + provenance", 7)


def slide_8_wrap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _paint_bg(s, DARK)
    _text(s, "Status  +  thanks",
          Inches(0.7), Inches(0.4), Inches(12), Inches(0.7),
          size=28, bold=True, color=WHITE)
    _line(s, Inches(0.7), Inches(1.15), Inches(12.6), Inches(1.15),
          color=CREAM, weight=2.5)

    _text(s, "Done", Inches(0.7), Inches(1.3), Inches(5.8), Inches(0.4),
          size=15, bold=True, color=TEAL)
    _bullets(s, [
        "UCI Grid Stability + tau-robustness margin",
        "Synthetic dataset + IEEE 519 compliance classifier",
        "Multi-horizon THD forecasting",
        "Anomaly / cyber-resilience detection (both datasets, compared)",
        "4-class disturbance prediction (99.70% accuracy)",
        "SHAP explainability layer",
        "Interactive dashboard (Module 7)",
        "Dataset provenance confirmed with the team",
    ], Inches(0.7), Inches(1.8), Inches(5.8), Inches(4.5),
       size=12, color=LIGHT, spacing=6)

    _text(s, "Still open", Inches(7.0), Inches(1.3), Inches(5.8), Inches(0.4),
          size=15, bold=True, color=CORAL)
    _bullets(s, [
        "Re-derive economic_cost_BDT from an explicit tariff + "
        "replacement-cost formula",
        "Binary breach classifier for forecasting's early-warning F1",
        "Multi-seed + bootstrap CI on forecasting numbers",
        "LaTeX manuscript draft",
    ], Inches(7.0), Inches(1.8), Inches(5.8), Inches(2.5),
       size=12, color=LIGHT, spacing=6)

    strip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.7), Inches(5.9),
                                Inches(12.1), Inches(1.0))
    strip.fill.solid(); strip.fill.fore_color.rgb = CORAL
    strip.line.color.rgb = CORAL
    tf = strip.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25); tf.margin_top = Inches(0.10)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Thanks — questions welcome"
    r.font.size = Pt(19); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r = p2.add_run(); r.text = "github.com/Fahmidaca/microgrid-dt"
    r.font.size = Pt(13); r.font.color.rgb = CREAM; r.font.name = "Consolas"

    _footer(s, "Status + thanks", 8)


# ============ DRIVER =====================================================
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    slide_1_title(prs)
    slide_2_scope(prs)
    slide_3_uci(prs)
    slide_4_synthetic_downstream(prs)
    slide_5_anomaly_comparison(prs)
    slide_6_prediction_xai(prs)
    slide_7_dashboard_provenance(prs)
    slide_8_wrap(prs)
    out = OUT_DIR / "cse_overall_progress.pptx"
    prs.save(out)
    print(f"[pptx-cse] wrote {out} ({out.stat().st_size/1024:.1f} KB, "
          f"{len(prs.slides)} slides)")


if __name__ == "__main__":
    build()
